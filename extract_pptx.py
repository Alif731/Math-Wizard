from pptx import Presentation

prs = Presentation("Change Schema.pptx")
print(f"Total slides: {len(prs.slides)}")

for i, slide in enumerate(prs.slides):
    print(f"\n--- Slide {i+1} ---")
    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
        notes = slide.notes_slide.notes_text_frame.text.strip()
        if notes:
            print(f"[Notes]: {notes}")
    
    for j, shape in enumerate(slide.shapes):
        shape_type = shape.shape_type
        text = ""
        if shape.has_text_frame:
            text = " / ".join(p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip())
        
        if shape.has_table:
            table_text = []
            for row in shape.table.rows:
                row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_text:
                    table_text.append(" | ".join(row_text))
            text = " [Table]: " + " // ".join(table_text)
            
        if text:
            print(f"Shape {j} ({shape_type}): {text}")
